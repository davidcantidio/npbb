"""Base reader utilities for PPTX slide iteration.

This module provides a minimal, reusable PPTX reader for extractor pipelines.
It intentionally focuses on deterministic slide traversal and text inventory.

Limitations:
- Extracts text from shape text frames and table cells only.
- Falls back to ZIP/XML text extraction when `python-pptx` is unavailable.
- Does not parse chart values, SmartArt semantics, or embedded media binaries.
- Does not perform OCR on images.
- Does not include speaker notes by default.
"""

from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path
import re
from typing import Any, Iterator

try:  # package execution (repo root at `npbb`)
    from etl.extract.pptx_xml_fallback import PptxXmlFallbackError, iter_slide_text_blocks
except ModuleNotFoundError:  # pragma: no cover - fallback for `npbb.*` imports
    from npbb.etl.extract.pptx_xml_fallback import PptxXmlFallbackError, iter_slide_text_blocks

try:  # package execution (repo root at `npbb`)
    from core.registry.location_ref import format_location
except ModuleNotFoundError:  # pragma: no cover - fallback for `npbb.*` imports
    from npbb.core.registry.location_ref import format_location


_SPACE_RE = re.compile(r"\s+")


class PptxReaderError(RuntimeError):
    """Raised when a PPTX file cannot be opened or parsed safely."""


class PptxDependencyMissingError(PptxReaderError):
    """Raised when `python-pptx` is not available in current environment."""


@dataclass(frozen=True)
class PptxSlide:
    """Structured slide payload returned by :func:`iter_slides`.

    Attributes:
        slide_number: One-based index of the slide in the deck.
        slide_id: Stable internal slide identifier exposed by python-pptx.
        layout_name: Name of the slide layout, when available.
        text_items: Ordered text fragments extracted from shapes/tables.
        lineage_location: Canonical lineage location string (`slide:<n>`).
    """

    slide_number: int
    slide_id: int
    layout_name: str | None
    text_items: list[str]
    lineage_location: str


def _normalize_text(value: Any) -> str:
    """Normalize arbitrary values into compact single-line text tokens."""
    if value is None:
        return ""
    text = str(value).strip()
    if not text:
        return ""
    return _SPACE_RE.sub(" ", text)


def _extract_table_texts(shape: Any) -> list[str]:
    """Return normalized table cell texts from one shape, preserving order."""
    if not getattr(shape, "has_table", False):
        return []

    out: list[str] = []
    table = shape.table
    for row in table.rows:
        for cell in row.cells:
            text = _normalize_text(cell.text)
            if text:
                out.append(text)
    return out


def _extract_shape_texts(shape: Any) -> list[str]:
    """Return normalized text items for one shape (text frame + table)."""
    out: list[str] = []

    if getattr(shape, "has_text_frame", False):
        for paragraph in shape.text_frame.paragraphs:
            text = _normalize_text(paragraph.text)
            if text:
                out.append(text)

    out.extend(_extract_table_texts(shape))
    return out


def _load_presentation(pptx_path: Path) -> Any:
    """Load a PPTX presentation with explicit, actionable failure messages.

    Args:
        pptx_path: Filesystem path to a `.pptx` file.

    Returns:
        A `pptx.Presentation` object.

    Raises:
        PptxReaderError: If dependency is missing, file does not exist, or file
            is not a valid PPTX package.
    """
    try:
        from pptx import Presentation  # type: ignore import-not-found
    except ModuleNotFoundError as exc:  # pragma: no cover - environment dependent
        raise PptxDependencyMissingError(
            "Dependencia ausente: instale 'python-pptx' para habilitar o extractor PPTX."
        ) from exc

    resolved = pptx_path.expanduser()
    if not resolved.exists():
        raise PptxReaderError(f"Arquivo PPTX nao encontrado: {resolved}")
    if not resolved.is_file():
        raise PptxReaderError(f"Caminho invalido para PPTX (nao e arquivo): {resolved}")

    try:
        return Presentation(str(resolved))
    except Exception as exc:  # pragma: no cover - depends on parser internals
        raise PptxReaderError(
            f"Falha ao abrir PPTX '{resolved}'. Verifique se o arquivo e um .pptx valido."
        ) from exc


def _slide_location(slide_number: int) -> str:
    """Return canonical lineage location for one slide number."""
    try:
        return format_location("slide", {"number": slide_number})
    except Exception:  # pragma: no cover - defensive fallback
        return f"slide:{slide_number}"


def _iter_slides_python_pptx(pptx_path: Path) -> Iterator[PptxSlide]:
    """Iterate slides using `python-pptx` parser."""
    presentation = _load_presentation(pptx_path)
    for index, slide in enumerate(presentation.slides, start=1):
        text_items: list[str] = []
        for shape in slide.shapes:
            text_items.extend(_extract_shape_texts(shape))
        yield PptxSlide(
            slide_number=index,
            slide_id=int(getattr(slide, "slide_id", index)),
            layout_name=getattr(getattr(slide, "slide_layout", None), "name", None),
            text_items=text_items,
            lineage_location=_slide_location(index),
        )


def _iter_slides_xml_fallback(pptx_path: Path) -> Iterator[PptxSlide]:
    """Iterate slides using ZIP/XML fallback parser."""
    try:
        for slide in iter_slide_text_blocks(pptx_path):
            text_items = list(slide.texts)
            if slide.slide_title and slide.slide_title not in text_items:
                text_items.insert(0, slide.slide_title)
            yield PptxSlide(
                slide_number=slide.slide_number,
                slide_id=slide.slide_number,
                layout_name=None,
                text_items=text_items,
                lineage_location=slide.lineage_location,
            )
    except PptxXmlFallbackError as exc:
        raise PptxReaderError(str(exc)) from exc


def iter_slides(
    pptx_path: str | Path,
    *,
    prefer_xml_fallback: bool = False,
) -> Iterator[PptxSlide]:
    """Iterate slides from a PPTX and return normalized text inventory.

    Args:
        pptx_path: Filesystem path to source presentation.
        prefer_xml_fallback: When `True`, skip `python-pptx` and use ZIP/XML
            fallback directly.

    Yields:
        :class:`PptxSlide` items in deck order, including text extracted from
        shape text frames/tables or fallback XML blocks.

    Raises:
        PptxReaderError: If the source file is missing, invalid, unreadable, or
            cannot be parsed by both primary and fallback readers.
    """
    resolved = Path(pptx_path).expanduser()
    if prefer_xml_fallback:
        yield from _iter_slides_xml_fallback(resolved)
        return

    try:
        yield from _iter_slides_python_pptx(resolved)
    except PptxDependencyMissingError:
        yield from _iter_slides_xml_fallback(resolved)
